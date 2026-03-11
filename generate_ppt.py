from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Define content (Updated with Project Explanation Mode details)
    slides_content = [
        {
            "title": "Defake: Advanced Deepfake Detection System",
            "subtitle": "Combating Misinformation through AI-Powered Video Analysis",
            "layout": 0 # Title Slide
        },
        {
            "title": "Problem Statement",
            "content": [
                "The Rise of Deepfakes: AI-generated synthetic media (e.g., DeepFaceLab, FaceSwap).",
                "Threats:",
                "  - Political misinformation and fake news.",
                "  - Financial fraud and identity theft.",
                "  - Reputation damage (non-consensual content).",
                "Need: A precise, automated forensic tool to verify video authenticity."
            ],
            "layout": 1
        },
        {
            "title": "Objectives",
            "content": [
                "1. Develop a Hybrid DL model (CNN + RNN) for spatial-temporal analysis.",
                "2. Implement Attention Mechanism to focus on manipulated frames.",
                "3. Create a full-stack web application (React + Flask).",
                "4. Achieve high accuracy using Transfer Learning (XceptionNet)."
            ],
            "layout": 1
        },
        {
            "title": "System Architecture",
            "content": [
                "Frontend (Client):",
                "  - React.js (v19) + Vite + TailwindCSS (Cyberpunk UI).",
                "  - Handles uploads and visualizes confidence scores.",
                "Backend (Server):",
                "  - Flask (Python) REST API.",
                "  - Handles file processing and model inference.",
                "AI Engine:",
                "  - TensorFlow/Keras.",
                "  - Pipeline: Video -> Preprocess -> Model -> Score."
            ],
            "layout": 1
        },
        {
            "title": "Deep Learning Model Architecture",
            "content": [
                "Type: Time-Distributed Convolutional Recurrent Neural Network",
                "1. Input: (Batch, 20 Frames, 299, 299, 3).",
                "2. Spatial Features: XceptionNet (ImageNet weights, TimeDistributed).",
                "3. Temporal Analysis: Bidirectional LSTM (256 units).",
                "4. Attention Mechanism: Custom layer to weight frame importance.",
                "5. Classification: Dense(128) -> Dropout(0.5) -> Sigmoid Output."
            ],
            "layout": 1
        },
        {
            "title": "Preprocessing Pipeline",
            "content": [
                "1. Frame Extraction: Extracts 20 uniformly distributed frames.",
                "2. Face Detection: MTCNN (Multi-task Cascaded CNN) detects and crops faces.",
                "3. Normalization: Resizing to 299x299, pixel values 0-1.",
                "4. Padding: Zero-padding for short videos.",
                "Note: Ensures model focuses on facial artifacts, not background."
            ],
            "layout": 1
        },
         {
            "title": "Technical Stack",
            "content": [
                "Frontend:",
                "  - React.js, Vite, Axios, Framer Motion.",
                "Backend & ML:",
                "  - Flask (API), TensorFlow/Keras (Model).",
                "  - OpenCV (Video Processing), NumPy (Matrix Ops).",
                " Infrastucture:",
                "  - Local GPU Support (CUDA)."
            ],
            "layout": 1
        },
        {
            "title": "Implementation Workflow",
            "content": [
                "1. User uploads video via React UI.",
                "2. POST /predict sends file to Flask.",
                "3. Backend saves temp file -> extracts frames -> detects faces.",
                "4. Model predicts score (0.0 - 1.0).",
                "5. JSON response: { label: 'FAKE', confidence: 0.98 }.",
                "6. UI updates with Red/Green indicator."
            ],
            "layout": 1
        },
        {
            "title": "Challenges Solved",
            "content": [
                "Temporal Inconsistency (Flickering):",
                "  - Solved by Bidirectional LSTM.",
                "Variable Video Lengths:",
                "  - Solved by Uniform Frame Sampling (20 frames).",
                "Focus on Manipulation:",
                "  - Solved by Attention Mechanism."
            ],
            "layout": 1
        },
        {
            "title": "Future Enhancements",
            "content": [
                "- Audio-Visual Multimodal Detection.",
                "- Real-time Streaming (RTMP/WebRTC).",
                "- Explainable AI (Grad-CAM heatmaps).",
                "- Mobile/Edge Optimization (TFLite)."
            ],
            "layout": 1
        },
        {
            "title": "Conclusion",
            "content": [
                "Defake combines advanced Computer Vision with Sequence Modeling.",
                "Provides a robust, user-friendly tool for digital forensics.",
                "Effective mitigation against the threat of deepfakes."
            ],
            "layout": 1
        },
        {
            "title": "Q&A",
            "content": [
                "Thank You!",
                "Questions?"
            ],
            "layout": 1
        }
    ]

    for slide_info in slides_content:
        layout_index = slide_info.get("layout", 1)
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title = slide.shapes.title
        title.text = slide_info.get("title", "")

        # Content
        if layout_index == 0: # Title Slide
            subtitle = slide.placeholders[1]
            subtitle.text = slide_info.get("subtitle", "")
        else: # Content Slide
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.word_wrap = True
            
            content_list = slide_info.get("content", [])
            for i, line in enumerate(content_list):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Check for indentation
                if line.strip().startswith("-"):
                    p.text = line.strip().lstrip("- ").strip()
                    p.level = 1
                elif line.strip().startswith("1.") or line.strip().startswith("2.") or line.strip().startswith("3.") or line.strip().startswith("4.") or line.strip().startswith("5."):
                     p.text = line
                     p.level = 0
                elif line.strip().startswith("Type:") or line.strip().startswith("Infrastucture:") or line.strip().startswith("Note:"):
                     p.text = line
                     p.level = 0
                else:
                    p.text = line
                    p.level = 0

    output_path = "Defake_Project_Presentation_Detailed.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
