from pptx import Presentation
from pptx.util import Inches, Pt

def create_full_presentation():
    # Use the default template
    prs = Presentation()

    slides = [
        # Slide 1
        {
            "title": "Deepfake Video Detection System",
            "content": ["Project Presentation"],
            "layout": 0
        },
        # Slide 2
        {
            "title": "Project Overview",
            "content": [
                "Objective: Automate the process of fake video detection using state-of-the-art Neural Networks.",
                "Scope: The project covers data preprocessing, model training, a backend for serving predictions, and a modern frontend interface for user interaction.",
                "Outcome: A web application where users can upload videos and receive an instant classification of REAL or FAKE along with a confidence score."
            ],
            "layout": 1
        },
        # Slide 3
        {
            "title": "Problem Statement",
            "content": [
                "With the rapid advancement of Generative AI, creating hyper-realistic fake videos (deepfakes) has become easily accessible.",
                "Deepfakes can be used maliciously for misinformation, identity theft, political manipulation, and financial fraud.",
                "Human eyes can no longer reliably distinguish between authentic and manipulated media.",
                "There is a critical need for an automated, highly accurate tool to analyze and detect these manipulated temporal and spatial artifacts in videos."
            ],
            "layout": 1
        },
        # Slide 4
        {
            "title": "Why This Project is Needed",
            "content": [
                "Security & Trust: To restore trust in digital media by providing a reliable verification mechanism.",
                "Scalability: Manual verification is impossible at the scale of modern social media; automated deep learning models can process videos in seconds.",
                "Protection: To protect individuals and public figures from digital impersonation.",
                "Accessibility: Providing a clear, user-friendly interface allows non-technical users to verify video authenticity easily."
            ],
            "layout": 1
        },
        # Slide 5
        {
            "title": "Frontend Explanation",
            "content": [
                "Technologies Used:",
                "- React 19 (via Vite): For building the user interface.",
                "- Tailwind CSS: For utility-first, modern, and responsive styling.",
                "- Framer Motion: For smooth and dynamic animations.",
                "- React Dropzone: For handling drag-and-drop video uploads effortlessly.",
                "- Axios: For making HTTP requests to the backend API.",
                "- Lucide React: For scalable and beautiful icons."
            ],
            "layout": 1
        },
        # Slide 6
        {
            "title": "Why Those Frontend Technologies Were Chosen",
            "content": [
                "- React & Vite: Fast development server (Vite) and component-based architecture (React) allow for a highly interactive and modular UI.",
                "- Tailwind CSS: Rapid UI development. It allows for a premium aesthetic (dark modes, gradients, glassmorphism) without leaving the HTML.",
                "- Framer Motion: Adds a premium feel with micro-animations that make the application feel responsive and alive.",
                "- Axios: Simplifies handling async API requests and file uploads (FormData)."
            ],
            "layout": 1
        },
        # Slide 7
        {
            "title": "How the Frontend Works",
            "content": [
                "- The application loads a single-page interface with a futuristic, dark-themed aesthetic.",
                "- The main component is an upload zone (powered by React Dropzone) that accepts .mp4 and .avi files.",
                "- Once a file is dropped, the frontend displays file metadata, a 'processing' state with animations.",
                "- It sends the file via an Axios POST request using FormData to the backend /predict endpoint.",
                "- Upon receiving the response, it transitions to a Results view, displaying 'REAL' or 'FAKE' and a confidence meter."
            ],
            "layout": 1
        },
        # Slide 8
        {
            "title": "User Interface Flow",
            "content": [
                "1. Landing / Upload: User opens the app. A sleek drag-and-drop zone prompts the user.",
                "2. File Selection: User drops a video file. The UI validates the file and shows its details.",
                "3. Processing: User clicks 'Analyze Video'. The UI shows a loading state with spinner.",
                "4. Backend API Call: Video is sent to the backend endpoint.",
                "5. Result Display: The system reveals the prediction (Real/Fake) along with a percentage confidence bar.",
                "6. Reset: User can click 'Upload Another Video' to test again."
            ],
            "layout": 1
        },
        # Slide 9
        {
            "title": "Backend Detailed Explanation Overview",
            "content": [
                "The backend is built in Python and is responsible for two major pipelines:",
                "",
                "1. Model Training Pipeline: Scripts to preprocess data, extract faces, and train the deep learning model.",
                "2. Inference / API Pipeline: A Flask server that loads the trained model, receives user videos, processes them, and returns predictions."
            ],
            "layout": 1
        },
        # Slide 10
        {
            "title": "Backend File: app.py",
            "content": [
                "Purpose: The main entry point for the Flask web server.",
                "How it works:",
                "- Initializes a Flask app and configures Cross-Origin Resource Sharing (CORS).",
                "- Defines a /predict POST route that receives the video from the frontend, saves it temporarily, and calls model_utils.predict_video().",
                "- Returns the prediction result (Real/Fake, Confidence) as a JSON response.",
                "Why it is necessary: Acts as the bridge between the user's web browser (frontend) and the heavy deep learning inference code."
            ],
            "layout": 1
        },
        # Slide 11
        {
            "title": "Backend File: config.py",
            "content": [
                "Purpose: Centralizes all configuration variables, paths, and hyperparameters.",
                "How it works:",
                "- Defines absolute directory paths for the real dataset, fake dataset, and processed .npy data.",
                "- Sets critical model hyperparameters: IMG_SIZE = (299, 299), FRAME_COUNT = 20, BATCH_SIZE = 2, EPOCHS, and LEARNING_RATE.",
                "Why it is necessary: Keeps the codebase clean. Modifying hyperparameters or file paths requires a single change in this file instead of hunting through scripts."
            ],
            "layout": 1
        },
        # Slide 12
        {
            "title": "Backend File: debug_files.py & debug_mtcnn.py",
            "content": [
                "Purpose: Troubleshooting and environment validation scripts.",
                "How it works:",
                "- debug_files.py: Checks if the fake video directories exist and counts the .mp4 files using glob.",
                "- debug_mtcnn.py: Attempts to import the MTCNN library and initialize the face detector to ensure the environment is set up correctly.",
                "Why it is necessary: Crucial for debugging paths and dependency issues on new environments before running long training scripts."
            ],
            "layout": 1
        },
        # Slide 13
        {
            "title": "Backend File: model_utils.py",
            "content": [
                "Purpose: Contains the core Deep Learning architecture and inference logic.",
                "How it works:",
                "- build_model(): Constructs the complex model using EfficientNetB0, TimeDistributed wrappers, Spatial Dropout, Bidirectional LSTM, and Custom Attention Layer.",
                "- predict_video(): Takes a raw video path, calls preprocessing, runs the tensor through the loaded model, and calculates the confidence score.",
                "Why it is necessary: It houses the actual 'brain' of the project. Without this file, the complex neural network structure would not exist."
            ],
            "layout": 1
        },
        # Slide 14
        {
            "title": "Backend File: preprocess.py",
            "content": [
                "Purpose: Real-time video processing and face extraction.",
                "How it works:",
                "- extract_multiple_clips(): Divides a video into multiple segments and extracts a uniform sequence of frames (e.g., 20 frames per clip).",
                "- detect_and_crop_faces(): Scans frames, uses OpenCV Haar Cascades to detect the largest face, crops it with a margin, and normalizes pixel values to [0, 1].",
                "Why it is necessary: Neural networks require consistent data. This ensures every video is formatted as a sequence of 20 aligned, cropped faces."
            ],
            "layout": 1
        },
        # Slide 15
        {
            "title": "Backend File: preprocess_dataset.py",
            "content": [
                "Purpose: Offline data preparation for training.",
                "How it works:",
                "- Scans the massive dataset of REAL and FAKE videos defined in config.py.",
                "- Uses preprocess.py to extract multiple facial clips from every video.",
                "- Saves these preprocessed 4D NumPy arrays to the disk as .npy files.",
                "Why it is necessary: Preprocessing during training bottlenecks the GPU. Doing it offline and saving .npy arrays vastly speeds up training loops."
            ],
            "layout": 1
        },
        # Slide 16
        {
            "title": "Backend File: train.py",
            "content": [
                "Purpose: The main script mapping data generation and training execution.",
                "How it works:",
                "- Uses an AugmentedDataGenerator to dynamically load batches into memory, applying runtime data augmentation (flipping, brightness, noise).",
                "- Compiles the model using Callbacks like ModelCheckpoint, ReduceLROnPlateau, and EarlyStopping.",
                "- Fits the model and plots validation accuracy/loss charts.",
                "Why it is necessary: Responsible for finding the optimal weights for fake detection, generating the model_weights.h5 file."
            ],
            "layout": 1
        },
        # Slide 17
        {
            "title": "Backend File: test_generator.py",
            "content": [
                "Purpose: Validates the custom Data Generator.",
                "How it works:",
                "- Instantiates the VideoDataGenerator from train.py using a tiny subset of the preprocessed dataset.",
                "- Attempts to fetch a sample batch to verify shapes (e.g., (Batch Size, 20, 299, 299, 3)).",
                "Why it is necessary: Confirms that dynamic data loading works flawlessly before committing to a 10-hour, large-scale training execution."
            ],
            "layout": 1
        },
        # Slide 18
        {
            "title": "Technologies Used (Backend & ML)",
            "content": [
                "- Python: The core programming language.",
                "- Flask: A lightweight, fast web framework.",
                "- TensorFlow / Keras: Framework used to build and train the neural networks.",
                "- OpenCV (cv2): Computer vision library used for reading video frames and detecting faces via Haar Cascades.",
                "- NumPy: Handles N-dimensional arrays smoothly."
            ],
            "layout": 1
        },
        # Slide 19
        {
            "title": "Deep Learning Architecture Selection",
            "content": [
                "- EfficientNetB0 (Base): Chosen for optimal accuracy-efficiency trade-off; acts as our spatial feature extractor.",
                "- TimeDistributed Layer: Allows EfficientNet to process each frame independently across time.",
                "- Bidirectional LSTM: Analyzes temporal patterns forwards and backwards to detect unnatural movements or flickering.",
                "- Custom Attention Layer: Dynamically weights the importance of specific frames where deepfake artifacts are prominent."
            ],
            "layout": 1
        },
        # Slide 20
        {
            "title": "Project Working Flow (Input to Output)",
            "content": [
                "1. User Uploads video to frontend.",
                "2. Transmitted via HTTP to Flask Backend.",
                "3. Server extracts 20 distinct frames via OpenCV.",
                "4. Haar Cascades crops out faces, normalizing colors.",
                "5. EfficientNet + LSTM analyzes frame sequence.",
                "6. Network emits probability score [0.0 - 1.0].",
                "7. Frontend receives JSON Result and maps score to UI feedback."
            ],
            "layout": 1
        },
        # Slide 21
        {
            "title": "Results Achieved",
            "content": [
                "- Output Accuracy: The system provides robust binary classification (Real/Fake) distinguishing high-quality deepfakes from genuine videos.",
                "- Precise Confidence Scores: Accurately reflects algorithmic uncertainty (e.g., 94.5%).",
                "- Speed Optimization: By migrating to EfficientNet and Haar Cascades, inference is achieved in merely a few seconds on arbitrary hardware."
            ],
            "layout": 1
        },
        # Slide 22
        {
            "title": "Visual Output of the Project",
            "content": [
                "During the execution, the frontend establishes clear states:",
                "- Seamless Cyberpunk-Themed UI Dropzone.",
                "- Animated Processing Sequence informing the user.",
                "- Definitive Output Screen (Color-coded Red for Fake, Green for Real) rendering a visual progress bar indicating AI Confidence.",
                "- Robust server logs capturing prediction mechanics and confidence metric precision."
            ],
            "layout": 1
        },
        # Slide 23
        {
            "title": "Challenges Faced",
            "content": [
                "- Data Scarcity & Imbalance: Acquiring high-fidelity diverse deepfake variations against matched real variants.",
                "- System Resource Constraint: Holding simultaneous sequential 20-frame high-resolution sets caused severe memory exhaustion.",
                "- Preprocessing Bottleneck: Initial MTCNN architectures stalled inference times to minutes.",
                "- Acute Overfitting: Complex RNN parameters effectively memorized training sets but failed robust generalization."
            ],
            "layout": 1
        },
        # Slide 24
        {
            "title": "Solutions Implemented",
            "content": [
                "- Offline Processing: Persisting data as .npy blobs to permanently sidestep realtime RAM saturation during epochs.",
                "- Haar Cascade Paradigm: Deprecated MTCNN for optimized OpenCV detectors achieving 10x frame discovery acceleration.",
                "- Advanced Regularization Framework: Stacked SpatialDropout1D modules, embedded severe L2 Kernels, applied 0.2 Label Smoothing.",
                "- Base Model Rotation: Substituted Xception for EfficientNetB0 for superior convergence paradigms."
            ],
            "layout": 1
        },
        # Slide 25
        {
            "title": "Conclusion",
            "content": [
                "- Successfully demonstrated the union of spatial vision feature extraction combined closely with recurrent sequential assessment algorithms.",
                "- Flask and React segregation ensures immediate microservice scalability capabilities for future integrations.",
                "- Final UI masks the sheer complexity of spatial inference architecture under a simplified intuitive presentation paradigm."
            ],
            "layout": 1
        },
        # Slide 26
        {
            "title": "Future Scope",
            "content": [
                "- Audio-Visual Synthesis Detection: Integration of Wav2Vec models paralleled with the visual pipeline.",
                "- Native Cloud Implementations: Dockerized hosting over AWS SageMaker for elastic computational resourcing targeting real-time execution.",
                "- Next Generation Architectures: Migration pathways toward complete Video Vision Transformers (ViViT).",
                "- Expanded Explainability Maps: Grad-CAM overlay projections demonstrating frame-wise artifact triggers."
            ],
            "layout": 1
        },
        # Slide 27
        {
            "title": "Thank You",
            "content": [
                "Thank you for your time.",
                "Are there any questions regarding the architecture, the training process, or the frontend implementation?"
            ],
            "layout": 1
        }
    ]

    for slide_data in slides:
        layout_id = slide_data['layout']
        slide_layout = prs.slide_layouts[layout_id]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        if title:
            title.text = slide_data['title']
            
        # Add content if applicable
        if layout_id == 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.word_wrap = True
            
            for index, point in enumerate(slide_data['content']):
                if index == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Check if it should be an indented bullet point
                if point.startswith("- "):
                    p.text = point[2:]
                    p.level = 1
                else:
                    p.text = point
                    p.level = 0
        elif layout_id == 0: # Title Slide
            subtitle = slide.placeholders[1]
            if subtitle and len(slide_data['content']) > 0:
                subtitle.text = slide_data['content'][0]

    prs.save("Deepfake_Project_Presentation_27_Slides.pptx")
    print("Successfully generated Deepfake_Project_Presentation_27_Slides.pptx with 27 slides!")

if __name__ == "__main__":
    create_full_presentation()
